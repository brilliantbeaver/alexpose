"""Build the URTC 2026 poster as an editable, vector-first PowerPoint slide.

All diagrams and charts are native PowerPoint shapes so that they remain editable
after import into Google Slides.  Numerical claims are taken from the real-mode
artifacts produced by notebooks 06 through 10.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).with_name("urtc2026_sjepa_gait_poster.pptx")

W, H = 48.0, 36.0

NAVY = "102A43"
NAVY_2 = "173F5F"
INK = "17324D"
TEAL = "087F8C"
CYAN = "2BB3C0"
SKY = "DDF3F5"
PALE = "F5F9FC"
WHITE = "FFFFFF"
MUTED = "587085"
LINE = "C9D8E3"
CORAL = "EE6A5B"
CORAL_PALE = "FCEDEA"
AMBER = "F2B84B"
AMBER_PALE = "FFF5D9"
GREEN = "3D8B67"
GREEN_PALE = "EAF5EF"
RED = "C94B55"
RED_PALE = "FCEAEC"
PURPLE = "6750A4"
PURPLE_PALE = "EFEBFA"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    if radius:
        kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(1.0)
    return shape


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, line, radius)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=22,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font="Arial",
    margin=0.04,
    italic=False,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=22, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    for text, attrs in runs:
        r = p.add_run()
        r.text = text
        r.font.name = attrs.get("font", "Arial")
        r.font.size = Pt(attrs.get("size", size))
        r.font.bold = attrs.get("bold", False)
        r.font.italic = attrs.get("italic", False)
        r.font.color.rgb = rgb(attrs.get("color", color))
    return box


def line(slide, x1, y1, x2, y2, color=INK, width=2.0, dash=None):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    if dash is not None:
        shp.line.dash_style = dash
    return shp


def circle(slide, cx, cy, d, fill=WHITE, line_color=INK, width=1.5):
    shp = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        cx - d / 2,
        cy - d / 2,
        d,
        d,
        fill,
        line_color,
    )
    shp.line.width = Pt(width)
    return shp


def pill(slide, text, x, y, w, h, fill=SKY, color=TEAL, size=17, bold=True, line_color=None):
    shp = add_rect(slide, x, y, w, h, fill, line_color or fill, radius=True)
    add_text(slide, text, x + 0.08, y + 0.01, w - 0.16, h - 0.02, size, color, bold, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shp


def card(slide, x, y, w, h, title, accent=TEAL, title_size=27):
    add_rect(slide, x, y, w, h, WHITE, LINE, radius=True)
    add_rect(slide, x, y, 0.18, h, accent, accent, radius=True)
    add_text(slide, title, x + 0.42, y + 0.18, w - 0.72, 0.55, title_size, INK, True)
    line(slide, x + 0.42, y + 0.86, x + w - 0.34, y + 0.86, LINE, 1.1)


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.3):
    shp = line(slide, x1, y1, x2, y2, color, width)
    shp.line.end_arrowhead = True
    return shp


def draw_camera(slide, x, y, s=1.0, color=TEAL):
    add_rect(slide, x, y + 0.18 * s, 0.78 * s, 0.52 * s, WHITE, color, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x + 0.63 * s, y + 0.19 * s, 0.42 * s, 0.48 * s, color, color)
    circle(slide, x + 0.35 * s, y + 0.44 * s, 0.22 * s, SKY, color, 1.2)


def draw_skeleton(slide, x, y, scale=1.0, highlight=True):
    """Draw a compact BlazePose-style skeleton using native lines and circles."""
    pts = {
        "head": (x + 2.0 * scale, y + 0.5 * scale),
        "neck": (x + 2.0 * scale, y + 1.35 * scale),
        "ls": (x + 1.32 * scale, y + 1.60 * scale),
        "rs": (x + 2.68 * scale, y + 1.60 * scale),
        "le": (x + 0.90 * scale, y + 2.55 * scale),
        "re": (x + 3.10 * scale, y + 2.55 * scale),
        "lw": (x + 0.70 * scale, y + 3.45 * scale),
        "rw": (x + 3.30 * scale, y + 3.45 * scale),
        "lh": (x + 1.55 * scale, y + 3.55 * scale),
        "rh": (x + 2.45 * scale, y + 3.55 * scale),
        "lk": (x + 1.35 * scale, y + 5.15 * scale),
        "rk": (x + 2.65 * scale, y + 5.15 * scale),
        "la": (x + 1.10 * scale, y + 6.75 * scale),
        "ra": (x + 2.90 * scale, y + 6.75 * scale),
        "lf": (x + 0.68 * scale, y + 7.05 * scale),
        "rf": (x + 3.32 * scale, y + 7.05 * scale),
    }
    edges = [
        ("head", "neck"), ("neck", "ls"), ("neck", "rs"), ("ls", "rs"),
        ("ls", "le"), ("le", "lw"), ("rs", "re"), ("re", "rw"),
        ("ls", "lh"), ("rs", "rh"), ("lh", "rh"),
        ("lh", "lk"), ("lk", "la"), ("la", "lf"),
        ("rh", "rk"), ("rk", "ra"), ("ra", "rf"),
    ]
    for a, b in edges:
        line(slide, *pts[a], *pts[b], NAVY_2, max(1.3, 2.3 * scale))
    circle(slide, *pts["head"], 0.64 * scale, WHITE, NAVY_2, 1.8)
    selected = {"ls", "rs", "lh", "rh", "lk", "rk", "la", "ra", "lf", "rf"}
    for name, (px, py) in pts.items():
        if name == "head":
            continue
        active = highlight and name in selected
        circle(
            slide,
            px,
            py,
            0.25 * scale if not active else 0.36 * scale,
            CORAL if active else WHITE,
            CORAL if active else NAVY_2,
            1.3,
        )


def draw_step_icon(slide, kind, x, y, color):
    if kind == "video":
        draw_camera(slide, x + 0.20, y + 0.08, 0.85, color)
    elif kind == "pose":
        draw_skeleton(slide, x - 0.05, y - 0.03, 0.15, True)
    elif kind == "tokens":
        for rr in range(4):
            for cc in range(5):
                fill = color if (rr + cc) % 3 == 0 else SKY
                add_rect(slide, x + 0.12 + cc * 0.19, y + 0.12 + rr * 0.19, 0.14, 0.14, fill, fill, radius=True)
    elif kind == "mask":
        for rr in range(4):
            for cc in range(5):
                fill = CORAL if (rr * 2 + cc) % 3 == 0 else PALE
                add_rect(slide, x + 0.12 + cc * 0.19, y + 0.12 + rr * 0.19, 0.14, 0.14, fill, fill, radius=True)
    elif kind == "latent":
        for i, hh in enumerate([0.30, 0.58, 0.42, 0.70, 0.50]):
            add_rect(slide, x + 0.16 + i * 0.18, y + 0.84 - hh, 0.11, hh, color, color, radius=True)
    elif kind == "forest":
        for dx, dy, rr in [(0.30, 0.28, 0.28), (0.62, 0.19, 0.33), (0.88, 0.30, 0.25)]:
            circle(slide, x + dx, y + dy, rr * 2, GREEN_PALE, GREEN, 1.3)
            line(slide, x + dx, y + dy + rr, x + dx, y + 0.92, GREEN, 2)


def horizontal_bars(slide, items, x, y, w, row_h, max_value=1.0, label_w=4.1, value_fmt="{:.1%}"):
    for i, (label, value, color) in enumerate(items):
        yy = y + i * row_h
        add_text(slide, label, x, yy + 0.03, label_w - 0.15, row_h - 0.05, 18, INK, value == max(v for _, v, _ in items), valign=MSO_ANCHOR.MIDDLE)
        bx = x + label_w
        bw = w - label_w - 1.05
        add_rect(slide, bx, yy + 0.18, bw, 0.34, "E8EFF4", "E8EFF4", radius=True)
        add_rect(slide, bx, yy + 0.18, max(0.04, bw * value / max_value), 0.34, color, color, radius=True)
        add_text(slide, value_fmt.format(value), x + w - 0.95, yy, 0.92, row_h, 18, color, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background and header.
    add_rect(slide, 0, 0, W, H, PALE, PALE)
    add_rect(slide, 0, 0, W, 5.35, NAVY, NAVY)
    add_rect(slide, 0, 5.35, W, 1.05, TEAL, TEAL)
    add_text(
        slide,
        "Learning Monocular Gait Representations through\nNeurologically Guided Skeleton JEPA",
        0.95, 0.42, 30.2, 2.35, 46, WHITE, True, font="Arial", line_spacing=0.88,
    )
    add_text(
        slide,
        "Alex Mui · Penny Inouye · Theodore Mui · Phil Mui",
        1.02, 2.86, 28.8, 0.50, 23, SKY, True,
    )
    add_text(
        slide,
        "Aspiring Scholars Directed Research Program (ASDRP) · Fremont, California, USA",
        1.02, 3.43, 28.8, 0.42, 18, WHITE,
    )
    pill(slide, "URTC 2026 · Research use only", 1.02, 4.15, 6.25, 0.58, CYAN, NAVY, 16, True)
    draw_skeleton(slide, 31.25, 0.43, 0.43, True)
    add_text(slide, "SOURCE-VIDEO AUDIT", 36.0, 0.62, 10.7, 0.46, 18, AMBER, True, PP_ALIGN.CENTER)
    add_text(slide, "65.6%", 35.9, 1.05, 4.4, 1.15, 42, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    arrow(slide, 40.22, 1.62, 41.40, 1.62, AMBER, 3.6)
    add_text(slide, "24.0%", 41.35, 1.05, 4.5, 1.15, 42, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, "6-fold CV · sequence IDs", 35.75, 2.10, 4.75, 0.35, 14, SKY, align=PP_ALIGN.CENTER)
    add_text(slide, "6-fold CV · source videos", 41.30, 2.10, 4.75, 0.35, 14, SKY, align=PP_ALIGN.CENTER)
    add_rect(slide, 36.0, 2.75, 9.9, 1.45, "193B55", CYAN, radius=True)
    add_text(slide, "−41.7 percentage points", 36.2, 2.93, 9.5, 0.46, 25, AMBER, True, PP_ALIGN.CENTER)
    add_text(slide, "matched grouping gap", 36.2, 3.50, 9.5, 0.30, 15, WHITE, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Source-video holdout changes the conclusion: 24.0% accuracy, below the 49.0% majority control.",
        1.0, 5.56, 46.0, 0.55, 23, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
    )

    x1, x2, x3 = 0.75, 16.20, 31.65
    cw = 14.70

    # LEFT COLUMN: question.
    card(slide, x1, 6.75, cw, 4.05, "1 · Research question", CYAN)
    add_text(
        slide,
        "Can normal-only, label-free S-JEPA pretraining learn useful single-camera gait features, and do they survive a source-video holdout?",
        x1 + 0.48, 7.85, cw - 0.96, 1.30, 24, INK, True, line_spacing=0.95,
    )
    add_text(
        slide,
        "We trace the full path from RGB video to pose tokens, frozen embeddings, condition probes, source-identity audits, mask ablations, and causal latent forecasts.",
        x1 + 0.48, 9.20, cw - 0.96, 1.12, 19, MUTED, line_spacing=0.95,
    )

    # LEFT COLUMN: cohort table.
    card(slide, x1, 11.10, cw, 7.05, "2 · Small, source-concentrated cohort", CORAL)
    add_text(slide, "GAVD subset", x1 + 0.50, 12.12, 4.0, 0.38, 17, MUTED, True)
    add_text(slide, "96 sequences", x1 + 0.50, 12.48, 5.0, 0.64, 29, TEAL, True)
    add_text(slide, "18 source videos", x1 + 5.15, 12.48, 4.8, 0.64, 29, CORAL, True)
    add_text(slide, "Condition", x1 + 0.55, 13.40, 5.0, 0.36, 17, MUTED, True)
    add_text(slide, "Sequences", x1 + 8.7, 13.40, 2.4, 0.36, 17, MUTED, True, PP_ALIGN.RIGHT)
    add_text(slide, "Videos", x1 + 11.4, 13.40, 2.1, 0.36, 17, MUTED, True, PP_ALIGN.RIGHT)
    rows = [
        ("Normal", 12, 1, CORAL_PALE),
        ("Parkinson’s disease", 9, 2, WHITE),
        ("Stroke", 12, 3, WHITE),
        ("Cerebral palsy", 16, 2, WHITE),
        ("Myopathic gait", 47, 10, WHITE),
    ]
    yy = 13.88
    for label, nseq, nvid, fill in rows:
        add_rect(slide, x1 + 0.48, yy, cw - 0.96, 0.58, fill, fill, radius=True)
        add_text(slide, label, x1 + 0.62, yy + 0.03, 6.9, 0.48, 18, INK, label == "Normal", valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, str(nseq), x1 + 8.85, yy + 0.03, 2.0, 0.48, 18, INK, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        add_text(slide, str(nvid), x1 + 11.55, yy + 0.03, 1.8, 0.48, 18, CORAL if label == "Normal" else INK, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        yy += 0.65
    pill(slide, "All 12 normal clips come from one upload", x1 + 0.55, 17.15, 8.55, 0.55, CORAL_PALE, RED, 16, True)
    add_text(slide, "This prevents a valid video-disjoint five-class estimate.", x1 + 9.20, 17.19, 4.7, 0.48, 15, MUTED, True)

    # LEFT COLUMN: anatomic prior.
    card(slide, x1, 18.45, cw, 10.35, "3 · Literature guides where to mask", PURPLE)
    cues = [
        ("Parkinson’s", "cadence · step length · support time", PURPLE_PALE, PURPLE),
        ("Stroke", "laterality · temporal asymmetry", SKY, TEAL),
        ("Cerebral palsy", "flexed hip–knee–ankle chain", AMBER_PALE, "A56B00"),
        ("Myopathy", "proximal weakness · trunk compensation", GREEN_PALE, GREEN),
    ]
    cy = 19.65
    for name, detail, fill, color in cues:
        add_rect(slide, x1 + 0.55, cy, 7.75, 1.35, fill, fill, radius=True)
        add_text(slide, name, x1 + 0.78, cy + 0.14, 6.9, 0.34, 18, color, True)
        add_text(slide, detail, x1 + 0.78, cy + 0.57, 6.9, 0.52, 16, INK)
        cy += 1.55
    draw_skeleton(slide, x1 + 9.00, 19.55, 1.03, True)
    add_text(slide, "10 bilateral targets", x1 + 8.65, 27.05, 5.55, 0.38, 18, CORAL, True, PP_ALIGN.CENTER)
    add_text(slide, "shoulders · hips · knees\nankles · foot indices", x1 + 8.65, 27.44, 5.55, 0.85, 16, MUTED, align=PP_ALIGN.CENTER, line_spacing=0.90)

    # LEFT COLUMN: validity boundary.
    card(slide, x1, 29.10, cw, 4.75, "Interpretation boundary", RED)
    add_text(
        slide,
        "These are video-level gait labels, not diagnoses. MediaPipe depth is not calibrated 3-D motion capture. Results do not measure transfer to a new person, camera, or clinic.",
        x1 + 0.55, 30.28, cw - 1.10, 1.68, 20, INK, True, line_spacing=0.98,
    )
    pill(slide, "Feasibility study", x1 + 0.58, 32.35, 3.55, 0.58, RED_PALE, RED, 17, True)
    pill(slide, "Non-diagnostic", x1 + 4.35, 32.35, 3.55, 0.58, RED_PALE, RED, 17, True)
    pill(slide, "Transductive", x1 + 8.12, 32.35, 3.45, 0.58, RED_PALE, RED, 17, True)

    # MIDDLE COLUMN: pipeline.
    card(slide, x2, 6.75, cw, 7.10, "4 · End-to-end audited pipeline", TEAL)
    steps = [
        ("video", "RGB clips", "GAVD boxes"),
        ("pose", "33 landmarks", "MediaPipe"),
        ("tokens", "64 frames", "528 tokens"),
        ("mask", "Mask ≤60%", "eligible tokens"),
        ("latent", "384-d pool", "frozen encoder"),
        ("forest", "5 classes", "Random Forest"),
    ]
    sx = x2 + 0.45
    sw = 2.03
    gap = 0.23
    for i, (kind, label, detail) in enumerate(steps):
        xx = sx + i * (sw + gap)
        fill = CORAL_PALE if kind == "mask" else SKY if i < 4 else GREEN_PALE
        accent = CORAL if kind == "mask" else TEAL if i < 4 else GREEN
        add_rect(slide, xx, 8.12, sw, 3.38, fill, fill, radius=True)
        draw_step_icon(slide, kind, xx + 0.45, 8.38, accent)
        add_text(slide, label, xx + 0.10, 9.55, sw - 0.20, 0.55, 17, INK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(slide, detail, xx + 0.10, 10.17, sw - 0.20, 0.64, 14, MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            arrow(slide, xx + sw + 0.03, 9.82, xx + sw + gap - 0.04, 9.82, MUTED, 1.8)
    add_text(
        slide,
        "Pretraining uses 12 normal clips from one upload, no condition labels, 300 epochs, seed 42, and a frozen EMA target encoder.",
        x2 + 0.55, 12.05, cw - 1.10, 0.88, 18, MUTED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
    )

    # MIDDLE COLUMN: JEPA mechanism.
    card(slide, x2, 14.15, cw, 8.25, "5 · Predict a hidden representation, not coordinates", CYAN)
    # Token grid.
    gx, gy = x2 + 0.65, 15.55
    for rr in range(5):
        for cc in range(10):
            eligible = rr >= 2
            masked = eligible and ((rr * 3 + cc * 2) % 5 < 3)
            fill = CORAL if masked else CYAN if eligible else "DCE6EE"
            add_rect(slide, gx + cc * 0.37, gy + rr * 0.37, 0.28, 0.28, fill, fill, radius=True)
    add_text(slide, "joint groups", gx - 0.05, gy + 1.95, 3.7, 0.30, 14, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "time →", gx + 3.72, gy + 0.63, 1.15, 0.35, 14, MUTED)
    add_text(slide, "visible", gx + 0.05, gy + 2.42, 1.05, 0.32, 15, TEAL, True)
    add_rect(slide, gx + 1.05, gy + 2.45, 0.28, 0.28, CYAN, CYAN, radius=True)
    add_text(slide, "target", gx + 1.55, gy + 2.42, 1.05, 0.32, 15, CORAL, True)
    add_rect(slide, gx + 2.50, gy + 2.45, 0.28, 0.28, CORAL, CORAL, radius=True)

    # Architecture paths.
    bx = x2 + 5.52
    add_rect(slide, bx, 15.52, 2.55, 1.05, SKY, TEAL, radius=True)
    add_text(slide, "View encoder", bx, 15.66, 2.55, 0.40, 19, TEAL, True, PP_ALIGN.CENTER)
    add_rect(slide, bx + 3.10, 15.52, 2.55, 1.05, AMBER_PALE, "A56B00", radius=True)
    add_text(slide, "Predictor", bx + 3.10, 15.66, 2.55, 0.40, 19, "8B6100", True, PP_ALIGN.CENTER)
    arrow(slide, bx + 2.62, 16.04, bx + 3.00, 16.04, TEAL, 2.2)
    add_rect(slide, bx + 3.10, 18.36, 2.55, 1.05, PURPLE_PALE, PURPLE, radius=True)
    add_text(slide, "EMA teacher", bx + 3.10, 18.50, 2.55, 0.40, 18, PURPLE, True, PP_ALIGN.CENTER)
    arrow(slide, bx + 4.38, 18.20, bx + 4.38, 16.72, PURPLE, 2.2)
    pill(slide, "no gradients", bx + 5.98, 18.58, 2.10, 0.48, PURPLE_PALE, PURPLE, 14, True)
    add_text(slide, "visible tokens", bx + 0.05, 16.74, 2.45, 0.34, 15, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "complete clip", bx + 3.15, 19.62, 2.45, 0.34, 15, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "Cross-entropy\nbetween student and\nteacher latents", bx + 5.90, 15.46, 2.55, 1.62, 16, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rich_text(
        slide,
        [("The view encoder sees context; the EMA teacher defines stable latent targets. ", {"bold": True, "color": INK}),
         ("Uniform target sampling avoids assuming that high motion is always clinically important.", {"color": MUTED})],
        x2 + 0.65, 20.82, cw - 1.30, 0.92, 18,
    )

    # MIDDLE COLUMN: exact comparison.
    card(slide, x2, 22.70, cw, 11.15, "6 · Exact 47/21 comparison", AMBER)
    add_text(slide, "Five-class test accuracy", x2 + 0.55, 23.78, 5.0, 0.40, 18, MUTED, True)
    add_text(
        slide,
        "Same 21 test IDs · macro F1: S-JEPA 61.3%, handcrafted RF 72.8%",
        x2 + 5.05, 23.78, 8.95, 0.40, 14, MUTED, True, PP_ALIGN.RIGHT,
    )
    horizontal_bars(
        slide,
        [
            ("82-feature RF", 0.7619, AMBER),
            ("Frozen S-JEPA", 0.6190, TEAL),
            ("Missingness only", 0.3333, CORAL),
            ("Majority class", 0.3333, MUTED),
        ],
        x2 + 0.55, 24.35, cw - 1.10, 1.05, 1.0, 4.15,
    )
    line(slide, x2 + 5.73, 24.30, x2 + 5.73, 28.50, LINE, 1.0)
    # Training health mini metrics.
    add_rect(slide, x2 + 0.58, 28.92, 4.05, 1.52, GREEN_PALE, GREEN_PALE, radius=True)
    add_text(slide, "12.54 → 0.57", x2 + 0.78, 29.11, 3.65, 0.45, 23, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "training loss", x2 + 0.78, 29.62, 3.65, 0.30, 15, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, x2 + 4.88, 28.92, 4.05, 1.52, SKY, SKY, radius=True)
    add_text(slide, "0.340 → 0.412", x2 + 5.08, 29.11, 3.65, 0.45, 23, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "feature standard deviation", x2 + 5.08, 29.62, 3.65, 0.30, 15, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, x2 + 9.18, 28.92, 4.05, 1.52, PURPLE_PALE, PURPLE_PALE, radius=True)
    add_text(slide, "0.636 → 0.535", x2 + 9.38, 29.11, 3.65, 0.45, 23, PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, "mean pair cosine", x2 + 9.38, 29.62, 3.65, 0.30, 15, MUTED, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "The learned representation is non-collapsed and useful for this split, but the comparison is between complete systems with different feature pipelines.",
        x2 + 0.62, 31.00, cw - 1.24, 1.13, 19, INK, True, line_spacing=0.97,
    )
    pill(slide, "+28.6 points over missingness", x2 + 0.65, 32.40, 5.90, 0.62, GREEN_PALE, GREEN, 17, True)
    pill(slide, "−14.3 points vs handcrafted", x2 + 6.80, 32.40, 6.20, 0.62, AMBER_PALE, "8B6100", 17, True)

    # RIGHT COLUMN: validation ladder.
    card(slide, x3, 6.75, cw, 9.55, "7 · Source grouping changes the conclusion", RED)
    ladder = [
        ("Fixed 70/30 sequence split", "62.1%", "video-confounded", TEAL, SKY),
        ("6-fold CV · sequence IDs", "65.6 ± 11.7%", "same fold machinery", PURPLE, PURPLE_PALE),
        ("6-fold CV · source videos", "24.0 ± 16.5%", "unseen uploads", RED, RED_PALE),
    ]
    ly = 7.98
    for i, (label, value, note, color, fill) in enumerate(ladder):
        add_rect(slide, x3 + 0.60, ly, cw - 1.20, 1.62, fill, fill, radius=True)
        add_text(slide, label, x3 + 0.88, ly + 0.20, 6.5, 0.46, 20, INK, True)
        add_text(slide, note, x3 + 0.88, ly + 0.83, 5.5, 0.36, 15, MUTED)
        add_text(slide, value, x3 + 7.25, ly + 0.25, 6.15, 0.82, 29 if i != 1 else 25, color, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        ly += 1.92
    add_rect(slide, x3 + 0.62, 13.90, cw - 1.24, 1.48, NAVY, NAVY, radius=True)
    add_text(slide, "41.7-point matched grouping gap", x3 + 0.85, 14.10, cw - 1.70, 0.47, 25, AMBER, True, PP_ALIGN.CENTER)
    add_text(slide, "Video-grouped majority accuracy: 49.0%", x3 + 0.85, 14.69, cw - 1.70, 0.31, 15, WHITE, align=PP_ALIGN.CENTER)

    # RIGHT COLUMN: representation audits.
    card(slide, x3, 16.60, cw, 8.05, "8 · What did the representation retain?", TEAL)
    audit_rows = [
        ("Source identity", "49.4 ± 12.1%", "chance 5.6% · untrained 51.7% · raw pose 46.0%", CORAL, CORAL_PALE),
        ("Best linear gait proxy", "R² = 0.086", "knee excursion; most canonical proxies near zero", TEAL, SKY),
        ("Latent phase clock", "R² = −0.13", "grouped frame probe · null result", PURPLE, PURPLE_PALE),
        ("Causal forecast", "7.0–8.6×", "normal 0.36 · abnormal 2.54–3.11 · source-confounded", AMBER, AMBER_PALE),
    ]
    ay = 17.76
    for title, metric, note, color, fill in audit_rows:
        add_rect(slide, x3 + 0.55, ay, cw - 1.10, 1.30, fill, fill, radius=True)
        circle(slide, x3 + 1.00, ay + 0.65, 0.40, color, color, 1.0)
        add_text(slide, title, x3 + 1.38, ay + 0.17, 4.55, 0.38, 19, INK, True)
        add_text(slide, note, x3 + 1.38, ay + 0.67, 7.7, 0.32, 14, MUTED)
        add_text(slide, metric, x3 + 9.50, ay + 0.26, 3.95, 0.62, 23, color, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        ay += 1.48
    add_text(
        slide,
        "Source identity was already strong in untrained and raw-pose baselines. Most canonical gait proxies were not linearly readable.",
        x3 + 0.62, 23.82, cw - 1.24, 0.60, 16, MUTED, True, PP_ALIGN.CENTER,
    )

    # RIGHT COLUMN: mask ablation.
    card(slide, x3, 24.95, cw, 4.95, "9 · One-seed mask ablation: no clear winner", PURPLE, 25)
    mini = [
        ("Neurologic-10", 0.5714, PURPLE),
        ("Random-10", 0.6190, TEAL),
        ("Motion-aware-10", 0.5714, AMBER),
        ("Full-body-33", 0.6190, GREEN),
    ]
    horizontal_bars(slide, mini, x3 + 0.55, 26.00, cw - 1.10, 0.72, 0.8, 4.20)
    add_text(slide, "Same 21 test clips. One prediction separates 57.1% from 61.9%, so no arm earns a superiority claim.", x3 + 0.60, 29.13, cw - 1.20, 0.40, 14, MUTED, True, PP_ALIGN.CENTER)

    # RIGHT COLUMN: conclusion.
    card(slide, x3, 30.20, cw, 3.65, "Take-home message", GREEN, 25)
    add_text(
        slide,
        "At this scale, the evaluation unit changes the conclusion more than mask choice. The next valid test needs multiple normal source videos, grouped splits, uncertainty, and matched-compute baselines.",
        x3 + 0.55, 31.28, cw - 1.10, 1.45, 20, INK, True, line_spacing=0.96,
    )
    pill(slide, "Audit the experimental unit first", x3 + 2.65, 32.95, 8.95, 0.58, GREEN_PALE, GREEN, 17, True)

    # Footer.
    add_rect(slide, 0, 34.25, W, 1.75, NAVY, NAVY)
    add_text(
        slide,
        "Key sources: Zanardi et al., Sci Rep 2021 · Lauzière et al., IJPMR 2014 · Pandey et al., IJO 2023 · Maulet et al., Neurology 2023 · Assran et al., CVPR 2023 · Abdelfattah & Alahi, ECCV 2024 · Gil et al., GAVD 2020 · Kapoor & Narayanan, Patterns 2023 · Roberts et al., Ecography 2017",
        0.75, 34.52, 38.0, 0.72, 11, WHITE, valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Code + reproducible artifacts\ngithub.com/brilliantbeaver/alexpose/tree/main/penny/gavd3",
        39.2, 34.53, 8.05, 0.78, 14, SKY, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE,
    )

    # Document properties and accessibility notes.
    prs.core_properties.title = "Learning Monocular Gait Representations through Neurologically Guided Skeleton JEPA"
    prs.core_properties.subject = "MIT URTC 2026 poster"
    prs.core_properties.author = "Alex Mui; Penny Inouye; Theodore Mui; Phil Mui"
    prs.core_properties.comments = "Vector-first poster generated from real notebook artifacts."

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
